#!/usr/bin/env python3
"""
build_rm_fcst_excel.py
======================

data/rm_fcst.json → data/RM_FCST_정리.xlsx (+ docs/data 동기화) 재생성.

주간업무보고 'RM FCST 정리' 버튼이 내려주는 정적 엑셀을 rm_fcst.json 기준으로
다시 만든다. 대상월 = 당월(현재 달). rm_fcst.json에 당월이 없으면 _months_covered
첫 월로 폴백.

시트 3개:
  - 사업장별 총괄: 23사업장 + 권역 4행([비발디]/[한국중부]/[아시아퍼시픽]/[한국남부])
  - 세그먼트별 상세: 사업장×(OTA/G-OTA/Inbound)
  - 메타정보

단위: RN=실, ADR=천원, 매출=백만원(VAT제외). FCST/예산 = FCST RN / 예산 RN × 100.
"""

import json
from datetime import datetime
from pathlib import Path

import openpyxl
from openpyxl.styles import Font, PatternFill, Alignment, Border, Side

REPO = Path(__file__).resolve().parent.parent
SRC = REPO / "data" / "rm_fcst.json"
OUT = REPO / "data" / "RM_FCST_정리.xlsx"
DOCS_OUT = REPO / "docs" / "data" / "RM_FCST_정리.xlsx"
OUT_SOS = REPO / "data" / "소사업_예상매출.xlsx"
DOCS_OUT_SOS = REPO / "docs" / "data" / "소사업_예상매출.xlsx"
# 세일즈마케팅 예상매출현황 PPT: 템플릿(객실-only, 마케팅팀 제공)에서 시작 → 소사업 가산.
TEMPLATE_PPT = REPO / "data" / "templates" / "세일즈마케팅_예상매출현황_template.pptx"
OUT_PPT = REPO / "data" / "세일즈마케팅_예상매출현황.pptx"
DOCS_OUT_PPT = REPO / "docs" / "data" / "세일즈마케팅_예상매출현황.pptx"

SEGS = ["OTA", "G-OTA", "Inbound"]
REGION_LABELS = {
    "비발디": "[비발디]",
    "한국중부": "[한국중부]",
    "아시아퍼시픽": "[아시아퍼시픽]",
    "한국남부": "[한국남부]",
}

HDR_FILL = PatternFill("solid", fgColor="2F2F2F")
HDR_FONT = Font(bold=True, color="FFFFFF", size=10)
TITLE_FONT = Font(bold=True, size=12)
BOLD = Font(bold=True, size=10)
THIN = Side(style="thin", color="D0D0D0")
BORDER = Border(left=THIN, right=THIN, top=THIN, bottom=THIN)
CENTER = Alignment(horizontal="center", vertical="center")
RIGHT = Alignment(horizontal="right", vertical="center")


def _ratio(fcst_rn, bud_rn):
    return round(fcst_rn / bud_rn * 100, 1) if bud_rn else 0


def _style_header(ws, row, ncol):
    for c in range(1, ncol + 1):
        cell = ws.cell(row=row, column=c)
        cell.fill = HDR_FILL
        cell.font = HDR_FONT
        cell.alignment = CENTER
        cell.border = BORDER


def build_sosaup_workbook(props, target, sos, sos_meta, d, series):
    """소사업 예상매출 독립 엑셀 생성 (data/ + docs/data 동기화).

    소사업 매출(백만, VAT제외) = FCST 객실수 × 비율(천원/실) / 1000.
    """
    wb = openpyxl.Workbook()

    # ── 시트1: 사업장별 소사업 ──
    ws = wb.active
    ws.title = "사업장별 소사업"
    ws["A1"] = f"소사업 예상매출 - {target}  (객실 외, GS, VAT제외 / 단위: 백만원)"
    ws["A1"].font = TITLE_FONT
    hdr = ["사업장", "FCST 객실수", "OTA", "G-OTA", "Inbound", "소사업 합계(백만)"]
    for i, h in enumerate(hdr, 1):
        ws.cell(row=2, column=i, value=h)
    _style_header(ws, 2, len(hdr))

    r = 3
    grand = {seg: 0.0 for seg in SEGS}
    grand_rn = 0
    for name in props:
        node = props[name].get(target)
        if not node or name not in sos:
            continue
        segd = node.get("segments", {})
        per = {}
        rn_sum = 0
        for seg in SEGS:
            rn = segd.get(seg, {}).get("rm_fcst_rn", 0) or 0
            per[seg] = rn * sos[name].get(seg, 0) / 1000.0  # 백만
            rn_sum += rn
            grand[seg] += per[seg]
        grand_rn += rn_sum
        vals = [name, rn_sum, round(per["OTA"], 1), round(per["G-OTA"], 1),
                round(per["Inbound"], 1), round(sum(per.values()), 1)]
        for i, v in enumerate(vals, 1):
            ws.cell(row=r, column=i, value=v)
        r += 1
    # 합계 행
    tot_vals = ["합계", grand_rn, round(grand["OTA"], 1), round(grand["G-OTA"], 1),
                round(grand["Inbound"], 1), round(sum(grand.values()), 1)]
    for i, v in enumerate(tot_vals, 1):
        ws.cell(row=r, column=i, value=v).font = BOLD

    # ── 시트2: 세그먼트별 상세 ──
    ws2 = wb.create_sheet("세그먼트별 상세")
    ws2["A1"] = f"소사업 세그먼트별 - {target}"
    ws2["A1"].font = TITLE_FONT
    hdr2 = ["사업장", "세그먼트", "FCST 객실수", "소사업 비율(천원/실)", "소사업 매출(백만)"]
    for i, h in enumerate(hdr2, 1):
        ws2.cell(row=2, column=i, value=h)
    _style_header(ws2, 2, len(hdr2))

    r = 3
    for name in props:
        node = props[name].get(target)
        if not node or name not in sos:
            continue
        segd = node.get("segments", {})
        for seg in SEGS:
            rn = segd.get(seg, {}).get("rm_fcst_rn", 0) or 0
            ratio = sos[name].get(seg, 0)
            row = [name, seg, rn, round(ratio, 1), round(rn * ratio / 1000.0, 1)]
            for i, v in enumerate(row, 1):
                ws2.cell(row=r, column=i, value=v)
            r += 1

    # ── 시트3: 월별 누적 (마감월=온북 / 당월=RM FCST) ──
    ws4 = wb.create_sheet("월별 누적")
    ws4["A1"] = f"소사업 월별 누적 - {target} (마감월=온북 실적, 당월=RM FCST / 백만, VAT제외)"
    ws4["A1"].font = TITLE_FONT
    hdr4 = ["월", "객실수 기준", "OTA", "G-OTA", "Inbound", "소사업(백만)", "누계(백만)"]
    for i, h in enumerate(hdr4, 1):
        ws4.cell(row=2, column=i, value=h)
    _style_header(ws4, 2, len(hdr4))
    r = 3
    run = 0.0
    for it in series:
        run += it["total"]
        vals = [f"{it['month']}월", "온북" if it["closed"] else "RM FCST",
                round(it["seg"]["OTA"], 1), round(it["seg"]["G-OTA"], 1),
                round(it["seg"]["Inbound"], 1), round(it["total"], 1), round(run, 1)]
        for i, v in enumerate(vals, 1):
            ws4.cell(row=r, column=i, value=v)
        r += 1
    cumv = ["누계", "",
            round(sum(it["seg"]["OTA"] for it in series), 1),
            round(sum(it["seg"]["G-OTA"] for it in series), 1),
            round(sum(it["seg"]["Inbound"] for it in series), 1),
            round(sum(it["total"] for it in series), 1), round(run, 1)]
    for i, v in enumerate(cumv, 1):
        ws4.cell(row=r, column=i, value=v).font = BOLD

    # ── 시트4: 메타정보 ──
    ws3 = wb.create_sheet("메타정보")
    meta = [
        ("대상월", target),
        ("소사업 기준", f"전년 {sos_meta.get('ref','')} 실당소사업 × FCST 객실수 (미운영 사업장은 {sos_meta.get('fallback','')} 대체)"),
        ("소사업 정의", "GS 객실 외 7유형(식음+부대+아쿠아+스포츠+골프+유통+기타), VAT제외"),
        ("세그먼트", "인바운드=Inbound / 국내OTA=OTA / 해외OTA=G-OTA"),
        ("소사업 소스", sos_meta.get("src", "")),
        ("RM FCST 소스", d.get("_source_pdf", "")),
        ("스냅샷일자", d.get("_snapshot_date", "")),
    ]
    for i, (k, v) in enumerate(meta, 1):
        ws3.cell(row=i, column=1, value=k).font = BOLD
        ws3.cell(row=i, column=2, value=v)

    # ── 서식 ──
    for col in range(2, 7):
        for rc in ws.iter_rows(min_row=3, min_col=col, max_col=col):
            for cell in rc:
                if isinstance(cell.value, (int, float)):
                    cell.alignment = RIGHT
                    cell.number_format = "#,##0" if col == 2 else "#,##0.0"
    for col in range(3, 6):
        for rc in ws2.iter_rows(min_row=3, min_col=col, max_col=col):
            for cell in rc:
                if isinstance(cell.value, (int, float)):
                    cell.alignment = RIGHT
                    cell.number_format = "#,##0" if col == 3 else "#,##0.0"
    ws.column_dimensions["A"].width = 16
    ws.column_dimensions["F"].width = 16
    ws2.column_dimensions["A"].width = 16
    ws2.column_dimensions["D"].width = 18
    ws2.column_dimensions["E"].width = 15
    for col in range(3, 8):
        for rc in ws4.iter_rows(min_row=3, min_col=col, max_col=col):
            for cell in rc:
                if isinstance(cell.value, (int, float)):
                    cell.alignment = RIGHT
                    cell.number_format = "#,##0.0"
    ws4.column_dimensions["A"].width = 7
    ws4.column_dimensions["B"].width = 11
    ws4.column_dimensions["F"].width = 13
    ws4.column_dimensions["G"].width = 13
    ws3.column_dimensions["A"].width = 14
    ws3.column_dimensions["B"].width = 60

    wb.save(OUT_SOS)
    wb.save(DOCS_OUT_SOS)
    return round(sum(grand.values()), 1)


OTB_JSON = REPO / "docs" / "data" / "otb_data.json"


def compute_monthly_sosaup(target, sos_all, props):
    """월별 소사업 누적 — 마감월=온북(rns_actual), 당월(미마감)=RM FCST.

    소사업_seg(월) = Σ사업장 객실수 × 전년비율(해당 월) / 1000  [백만, VAT제외]
    반환: (series, cum, cur)
      series = [{"month": m, "closed": bool, "seg": {OTA,G-OTA,Inbound}, "total": v}, ...]
      cum    = {seg: 누계백만},  cur = {seg: 당월백만}
    otb_data.json 없으면 당월(FCST)만 계산해 반환.
    """
    year, cur_m = int(target.split("-")[0]), int(target.split("-")[1])
    try:
        otb = json.loads(OTB_JSON.read_text(encoding="utf-8")).get("allMonths", {})
    except Exception:
        otb = {}
    series, cum, cur = [], {s: 0.0 for s in SEGS}, {s: 0.0 for s in SEGS}
    for m in range(1, cur_m + 1):
        ym = f"{year}-{m:02d}"
        ratios = sos_all.get(ym, {}).get("ratios", {})
        closed = m < cur_m
        bps = otb.get(str(m), {}).get("byPropertySegment", {})
        seg_tot = {s: 0.0 for s in SEGS}
        for seg in SEGS:
            for row in bps.get(seg, []):
                prop = row.get("name")
                if prop not in props:
                    continue
                rn = (row.get("rns_actual") or 0) if closed else (row.get("rm_fcst_rn") or 0)
                seg_tot[seg] += rn * ratios.get(prop, {}).get(seg, 0) / 1000.0
        series.append({"month": m, "closed": closed, "seg": seg_tot,
                       "total": sum(seg_tot.values())})
        for s in SEGS:
            cum[s] += seg_tot[s]
        if m == cur_m:
            cur = dict(seg_tot)
    return series, cum, cur


def build_sales_pptx(target, cur_seg, cum_seg):
    """세일즈마케팅 예상매출현황 PPT 재생성.

    템플릿(객실-only)에서 시작 → GS Forecast = 객실 + 소사업(총매출):
      당월 = 객실 + 당월 소사업(미마감=FCST 객실수 기준),
      누계 = 객실 + 누계 소사업(마감월=온북 / 당월=FCST, Σ월별).
    증감/달성률/증감률은 템플릿의 목표·전년(이미 총매출 기준) 대비 재계산.
    매번 템플릿에서 시작하므로 이중가산 없음. python-pptx 미설치 시 건너뜀.
    반환: (당월 GS 소사업, 누계 GS 소사업) 백만 또는 None.
    """
    if not TEMPLATE_PPT.exists():
        return None
    try:
        from pptx import Presentation
        from pptx.dml.color import RGBColor
    except ImportError:
        print("⚠ python-pptx 미설치 — 예상매출현황 PPT 건너뜀 (xlsx만 생성)")
        return None

    cur = {s: round(cur_seg.get(s, 0)) for s in SEGS}
    cum = {s: round(cum_seg.get(s, 0)) for s in SEGS}
    cur_gs, cum_gs = sum(cur.values()), sum(cum.values())
    # 행(GS=7, OTA=8, GOTA=9, Inbound=10) → (당월 가산, 누계 가산)
    add = {7: (cur_gs, cum_gs), 8: (cur["OTA"], cum["OTA"]),
           9: (cur["G-OTA"], cum["G-OTA"]), 10: (cur["Inbound"], cum["Inbound"])}
    # (목표, FC, 증감, 달성, 전년, 전년증감, 증감률) — [당월, 누계]
    BLK = [(2, 4, 6, 7, 8, 10, 11), (13, 15, 17, 18, 19, 21, 22)]

    prs = Presentation(str(TEMPLATE_PPT))
    tbl = [sh for sh in prs.slides[0].shapes if sh.has_table][0].table

    def num(ri, ci):
        t = tbl.cell(ri, ci).text.replace(",", "").replace("%", "").strip()
        try:
            return float(t)
        except ValueError:
            return None

    def put(ri, ci, val, pct=False):
        par = tbl.cell(ri, ci).text_frame.paragraphs[0]
        if not par.runs:
            par.add_run()
        run = par.runs[0]
        run.text = f"{val:,.0f}%" if pct else f"{val:,.0f}"
        run.font.color.rgb = RGBColor(0xFF, 0, 0) if val < 0 else RGBColor(0, 0, 0)
        for extra in par.runs[1:]:
            extra._r.getparent().remove(extra._r)

    for ri, adds in add.items():
        for bi, (c_tgt, c_fc, c_diff, c_ach, c_ly, c_lyd, c_lyr) in enumerate(BLK):
            fc0 = num(ri, c_fc)
            if fc0 is None:
                continue
            fc = fc0 + adds[bi]
            put(ri, c_fc, fc)
            tgt = num(ri, c_tgt)
            ly = num(ri, c_ly)
            if tgt is not None:
                put(ri, c_diff, fc - tgt)
                put(ri, c_ach, round(fc / tgt * 100) if tgt else 0, pct=True)
            if ly is not None:
                put(ri, c_lyd, fc - ly)
                put(ri, c_lyr, round((fc - ly) / ly * 100) if ly else 0, pct=True)

    prs.save(str(OUT_PPT))
    prs.save(str(DOCS_OUT_PPT))
    return cur_gs, cum_gs


def main():
    d = json.loads(SRC.read_text(encoding="utf-8"))
    props = d["properties"]
    regions = d.get("regions", {})
    covered = d.get("_months_covered", [])

    cur = datetime.now().strftime("%Y-%m")
    target = cur if cur in covered else (covered[0] if covered else cur)

    # 소사업(객실 외) 비율 로드 — 대상월 비율이 있으면 소사업 매출 컬럼 추가.
    # 비율(천원/실, VAT제외)은 scripts/build_sosaup_ratio.py 가 소사업 엑셀에서 추출.
    # 소사업 매출(백만) = FCST RN × 비율 / 1000.
    sos_path = REPO / "data" / "sosaup_ratio.json"
    sos, sos_meta, sos_all = {}, {}, {}
    if sos_path.exists():
        sj = json.loads(sos_path.read_text(encoding="utf-8"))
        sos_all = sj.get("by_month", {})
        blk = sos_all.get(target, {})
        sos = blk.get("ratios", {})
        sos_meta = {"ref": blk.get("ref_year_month", ""),
                    "fallback": blk.get("fallback_month", ""),
                    "src": sj.get("_source_excel", "")}
    has_sos = bool(sos)

    def sosaup_mil(name, node):
        """사업장 node의 세그먼트별 (FCST RN × 비율) 합 → 소사업 매출(백만, VAT제외)."""
        if not has_sos or name not in sos:
            return None
        segd = node.get("segments", {})
        total = sum((segd.get(seg, {}).get("rm_fcst_rn", 0) or 0) * sos[name].get(seg, 0)
                    for seg in SEGS)
        return round(total / 1000, 1)

    wb = openpyxl.Workbook()

    # ── 시트1: 사업장별 총괄 ──
    ws = wb.active
    ws.title = "사업장별 총괄"
    ws["A1"] = f"RM FCST - {target}"
    ws["A1"].font = TITLE_FONT
    hdr = ["사업장", "예산 RN", "예산 ADR", "예산 매출(백만)",
           "FCST RN", "FCST ADR", "FCST 매출(백만)", "FCST/예산"]
    if has_sos:
        hdr += ["소사업 매출(백만)", "FCST+소사업(백만)"]
    n1 = len(hdr)
    ws.append([])  # row2 placeholder via direct write
    for i, h in enumerate(hdr, 1):
        ws.cell(row=2, column=i, value=h)
    _style_header(ws, 2, n1)

    def total_row(label, node):
        return [label, node["rm_budget_rn"], node["rm_budget_adr"], node["rm_budget_rev_mil"],
                node["rm_fcst_rn"], node["rm_fcst_adr"], node["rm_fcst_rev_mil"],
                _ratio(node["rm_fcst_rn"], node["rm_budget_rn"])]

    r = 3
    sum_fcst = sum_sos = 0.0  # 사업장 합계용(소사업)
    for name in props:
        node = props[name].get(target)
        if not node:
            continue
        for i, v in enumerate(total_row(name, node), 1):
            ws.cell(row=r, column=i, value=v)
        if has_sos:
            smil = sosaup_mil(name, node)
            if smil is not None:
                frev = node["rm_fcst_rev_mil"]
                ws.cell(row=r, column=9, value=smil)
                ws.cell(row=r, column=10, value=round(frev + smil, 1))
                sum_fcst += frev
                sum_sos += smil
        r += 1
    # 사업장 합계 행
    if has_sos:
        ws.cell(row=r, column=1, value="합계(사업장)").font = BOLD
        for col, val in ((7, round(sum_fcst, 1)), (9, round(sum_sos, 1)),
                         (10, round(sum_fcst + sum_sos, 1))):
            ws.cell(row=r, column=col, value=val).font = BOLD
        r += 1
    r += 1  # blank separator
    for rkey, label in REGION_LABELS.items():
        node = regions.get(rkey, {}).get(target)
        if not node:
            continue
        for i, v in enumerate(total_row(label, node), 1):
            cell = ws.cell(row=r, column=i, value=v)
            cell.font = BOLD
        r += 1

    # ── 시트2: 세그먼트별 상세 ──
    ws2 = wb.create_sheet("세그먼트별 상세")
    ws2["A1"] = f"RM FCST 세그먼트별 - {target}"
    ws2["A1"].font = TITLE_FONT
    hdr2 = ["사업장", "세그먼트", "예산 RN", "예산 ADR", "예산 매출(백만)",
            "FCST RN", "FCST ADR", "FCST 매출(백만)", "FCST/예산", "차이 RN"]
    if has_sos:
        hdr2 += ["소사업 비율(천원/실)", "소사업 매출(백만)"]
    n2 = len(hdr2)
    for i, h in enumerate(hdr2, 1):
        ws2.cell(row=2, column=i, value=h)
    _style_header(ws2, 2, n2)

    r = 3
    for name in props:
        node = props[name].get(target)
        if not node:
            continue
        segd = node.get("segments", {})
        for seg in SEGS:
            s = segd.get(seg, {})
            brn = s.get("rm_budget_rn", 0)
            frn = s.get("rm_fcst_rn", 0)
            row = [name, seg, brn, s.get("rm_budget_adr", 0), s.get("rm_budget_rev_mil", 0),
                   frn, s.get("rm_fcst_adr", 0), s.get("rm_fcst_rev_mil", 0),
                   _ratio(frn, brn), frn - brn]
            for i, v in enumerate(row, 1):
                ws2.cell(row=r, column=i, value=v)
            if has_sos and name in sos:
                ratio = sos[name].get(seg, 0)            # 천원/실, VAT제외
                ws2.cell(row=r, column=11, value=round(ratio, 1))
                ws2.cell(row=r, column=12, value=round(frn * ratio / 1000, 1))  # 백만
            r += 1

    # ── 시트3: 메타정보 ──
    ws3 = wb.create_sheet("메타정보")
    meta = [
        ("소스 PDF", d.get("_source_pdf", "")),
        ("추출일시", d.get("_extracted_at", "")),
        ("스냅샷일자", d.get("_snapshot_date", "")),
        ("대상월", target),
        ("비고", "OTA/G-OTA/Inbound 세그먼트 Forecast / 단위: RN=실, ADR=천원, 매출=백만원, VAT제외"),
    ]
    if has_sos:
        meta += [
            ("소사업 기준", f"전년 {sos_meta.get('ref','')} 실당소사업 × FCST 객실수 (미운영 사업장은 {sos_meta.get('fallback','')} 대체)"),
            ("소사업 정의", "GS 객실 외 7유형(식음+부대+아쿠아+스포츠+골프+유통+기타), VAT제외"),
            ("소사업 소스", sos_meta.get("src", "")),
        ]
    for i, (k, v) in enumerate(meta, 1):
        ws3.cell(row=i, column=1, value=k).font = BOLD
        ws3.cell(row=i, column=2, value=v)

    # ── 서식: 숫자 우측정렬 + 열너비 ──
    # 소수1자리 컬럼: FCST/예산 + 소사업 관련. 총괄/세그 각각 has_sos 여부로 폭 결정.
    ws_ncol = 10 if has_sos else 8
    ws2_ncol = 12 if has_sos else 10
    ws_dec = {8, 9, 10} if has_sos else {8}        # 총괄: FCST/예산, 소사업매출, FCST+소사업
    ws2_dec = {9, 11, 12} if has_sos else {9}      # 세그: FCST/예산, 소사업비율, 소사업매출
    for sheet, ncol, dec in ((ws, ws_ncol, ws_dec), (ws2, ws2_ncol, ws2_dec)):
        for col in range(2, ncol + 1):
            for row_cells in sheet.iter_rows(min_row=3, min_col=col, max_col=col):
                for cell in row_cells:
                    if isinstance(cell.value, (int, float)):
                        cell.alignment = RIGHT
                        cell.number_format = "#,##0.0" if cell.column in dec else "#,##0"
        sheet.column_dimensions["A"].width = 16
    ws.column_dimensions["H"].width = 10
    ws2.column_dimensions["B"].width = 10
    if has_sos:
        ws.column_dimensions["I"].width = 15   # 소사업 매출(백만)
        ws.column_dimensions["J"].width = 17   # FCST+소사업(백만)
        ws2.column_dimensions["K"].width = 16  # 소사업 비율(천원/실)
        ws2.column_dimensions["L"].width = 15  # 소사업 매출(백만)
    ws3.column_dimensions["A"].width = 12
    ws3.column_dimensions["B"].width = 50

    wb.save(OUT)
    wb.save(DOCS_OUT)
    print(f"✓ RM_FCST_정리.xlsx 재생성 (대상월={target}, 사업장={sum(1 for n in props if props[n].get(target))})")
    print(f"  → {OUT}")
    print(f"  → {DOCS_OUT}")

    # 소사업 예상매출 독립 파일 (비율 데이터가 있을 때만)
    if has_sos:
        # 월별 누적: 마감월=온북(otb rns_actual) / 당월=RM FCST × 전년비율
        series, cum_seg, cur_seg = compute_monthly_sosaup(target, sos_all, props)
        total = build_sosaup_workbook(props, target, sos, sos_meta, d, series)
        print(f"✓ 소사업_예상매출.xlsx 생성 (대상월={target}, 당월={total:,.1f} / 누계={sum(cum_seg.values()):,.0f}백만 VAT제외)")
        print(f"  → {OUT_SOS}")
        print(f"  → {DOCS_OUT_SOS}")
        # 세일즈마케팅 예상매출현황 PPT (템플릿 + 소사업: 당월=FCST, 누계=온북+FCST)
        res = build_sales_pptx(target, cur_seg, cum_seg)
        if res is not None:
            cur_gs, cum_gs = res
            print(f"✓ 세일즈마케팅_예상매출현황.pptx 생성 (GS 소사업 당월+{cur_gs:,} / 누계+{cum_gs:,}백만 → 총매출)")
            print(f"  → {OUT_PPT}")
            print(f"  → {DOCS_OUT_PPT}")


if __name__ == "__main__":
    main()
